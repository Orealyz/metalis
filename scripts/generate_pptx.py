from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x00, 0x35, 0x80)
BLUE    = RGBColor(0x00, 0x78, 0xD4)
LBLUE   = RGBColor(0xE6, 0xF2, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE= RGBColor(0xF5, 0xF7, 0xFA)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY    = RGBColor(0x6B, 0x7A, 0x8D)
LGRAY   = RGBColor(0xE8, 0xEC, 0xF0)
GREEN   = RGBColor(0x10, 0x7C, 0x10)
LGREEN  = RGBColor(0xDC, 0xF5, 0xDC)
ORANGE  = RGBColor(0xD8, 0x3B, 0x01)
LORANGE = RGBColor(0xFF, 0xED, 0xE8)
PURPLE  = RGBColor(0x5C, 0x2D, 0x91)
LPURPLE = RGBColor(0xF0, 0xE6, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Primitives ────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, radius=0, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def badge(slide, label, x, y, bg, fg=WHITE):
    w = Inches(1.5)
    h = Inches(0.38)
    rect(slide, x, y, w, h, bg)
    txt(slide, label, x, y + Inches(0.04), w, h, 11, bold=True, color=fg, align=PP_ALIGN.CENTER)

def card(slide, x, y, w, h, title, body_lines, accent=BLUE):
    rect(slide, x, y, w, h, WHITE, line_color=LGRAY, line_width=0.5)
    rect(slide, x, y, w, Inches(0.06), accent)
    txt(slide, title, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.45),
        14, bold=True, color=NAVY)
    cy = y + Inches(0.65)
    for line in body_lines:
        txt(slide, line, x + Inches(0.2), cy, w - Inches(0.4), Inches(0.4), 12, color=GRAY)
        cy += Inches(0.38)

def slide_header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), W, Inches(0.05), BLUE)
    txt(slide, title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.75),
        30, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.5), Inches(0.82), Inches(10), Inches(0.38),
            14, color=RGBColor(0xA0, 0xBE, 0xDF), italic=True)
    rect(slide, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

def pill(slide, text, x, y, bg, fg=WHITE):
    w = Inches(1.9)
    h = Inches(0.5)
    rect(slide, x, y, w, h, bg)
    txt(slide, text, x + Inches(0.1), y + Inches(0.07), w - Inches(0.2), h,
        13, bold=True, color=fg, align=PP_ALIGN.CENTER)

# ── SLIDE 1 — Titre ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
rect(s, 0, 0, W, H, OFFWHITE)
rect(s, 0, 0, Inches(0.5), H, NAVY)
rect(s, Inches(0.5), 0, Inches(0.07), H, BLUE)

# Bloc titre central
rect(s, Inches(1.2), Inches(1.8), Inches(10.9), Inches(3.9), WHITE,
     line_color=LGRAY, line_width=0.5)
rect(s, Inches(1.2), Inches(1.8), Inches(0.12), Inches(3.9), NAVY)

txt(s, "METALIS", Inches(1.6), Inches(2.1), Inches(10.2), Inches(1.1),
    54, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
txt(s, "Infrastructure virtualisée sur Proxmox VE",
    Inches(1.6), Inches(3.1), Inches(10.2), Inches(0.7),
    22, color=BLUE, align=PP_ALIGN.LEFT)
txt(s, "MSPR — Ynov 2026",
    Inches(1.6), Inches(3.85), Inches(10.2), Inches(0.5),
    15, color=GRAY, align=PP_ALIGN.LEFT)

# Badges technos — 6 badges centrés dans 13.33"
# largeur badge 1.9" + gap 0.15" → total 12.15" → départ à 0.59"
bx = Inches(0.59)
for label, bg in [("Proxmox VE", NAVY), ("Active Directory", BLUE),
                   ("Odoo 17", GREEN), ("WooCommerce", PURPLE),
                   ("WireGuard", ORANGE), ("Grafana", RGBColor(0xE0, 0x5C, 0x00))]:
    pill(s, label, bx, Inches(5.9), bg)
    bx += Inches(2.05)

txt(s, "15 min  •  12 slides", Inches(1.6), Inches(6.7), Inches(10), Inches(0.4),
    12, color=GRAY)

# ── SLIDE 2 — Contexte ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Contexte", "METALIS — PME ~40 personnes, production 2×8")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

cards = [
    ("NAS saturé",       ["Fichiers CAO (SolidWorks)", "Panne avec perte données"], ORANGE),
    ("ERP lent",         ["Odoo lent en heure de charge", "Chevauchement pics prod/e-comm"], BLUE),
    ("Accès distants",   ["Commerciaux bloqués", "Prestataire CNC non cadré"], PURPLE),
    ("Pas de sauvegarde",["Copies USB ad hoc", "Non vérifiées"], NAVY),
]
cx = Inches(0.4)
for title, lines, accent in cards:
    card(s, cx, Inches(1.5), Inches(3.0), Inches(2.6), title, lines, accent)
    cx += Inches(3.15)

txt(s, "Question clé : « Que se passe-t-il si tout s'arrête un vendredi après-midi ? »",
    Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
    16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── SLIDE 3 — Objectifs ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Objectifs du projet")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

items = [
    (BLUE,   "Centraliser l'infrastructure sur un hyperviseur moderne"),
    (GREEN,  "Remplacer le NAS par un serveur de fichiers robuste"),
    (NAVY,   "Dimensionner l'ERP pour absorber les pics de charge"),
    (ORANGE, "VPN sécurisé pour commerciaux et prestataire CNC"),
    (PURPLE, "Supervision avec alertes — Prometheus · Grafana · Loki"),
    (BLUE,   "Documenter les procédures de continuité (PCA/PRA)"),
]
y = Inches(1.55)
for color, text in items:
    rect(s, Inches(0.5), y + Inches(0.07), Inches(0.35), Inches(0.35), color)
    txt(s, text, Inches(1.05), y, Inches(11.5), Inches(0.5), 17, color=DARK)
    y += Inches(0.7)

# ── SLIDE 4 — Choix Proxmox ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Choix de l'hyperviseur", "Proxmox VE vs VirtualBox")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

# VirtualBox — écarté
rect(s, Inches(0.4), Inches(1.4), Inches(5.9), Inches(5.6), WHITE, line_color=LGRAY, line_width=0.5)
rect(s, Inches(0.4), Inches(1.4), Inches(5.9), Inches(0.06), ORANGE)
txt(s, "✗  VirtualBox — écarté", Inches(0.6), Inches(1.55), Inches(5.5), Inches(0.55),
    17, bold=True, color=ORANGE)
y = Inches(2.2)
for line in ["Type 2 — s'installe sur un OS hôte",
             "Risque instabilité (MAJ Windows)",
             "Non adapté à un usage 24/7",
             "Pas de snapshot planifié natif"]:
    txt(s, "—  " + line, Inches(0.7), y, Inches(5.4), Inches(0.5), 15, color=GRAY)
    y += Inches(0.62)

# Proxmox — retenu
rect(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.6), WHITE, line_color=LGRAY, line_width=0.5)
rect(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(0.06), GREEN)
txt(s, "✓  Proxmox VE — retenu", Inches(7.1), Inches(1.55), Inches(5.6), Inches(0.55),
    17, bold=True, color=GREEN)
y = Inches(2.2)
for line in ["Type 1 — bare-metal, pas d'OS hôte",
             "KVM + LXC natif — performances optimales",
             "Interface web centralisée",
             "Snapshots planifiables via GUI/cron",
             "Gratuit (sans support commercial)"]:
    txt(s, "✓  " + line, Inches(7.1), y, Inches(5.6), Inches(0.5), 15, color=DARK)
    y += Inches(0.62)

# ── SLIDE 5 — Architecture VMs ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Architecture — Inventaire")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

vms = [
    ("100", "ct-vpn",        "192.168.1.208", "WireGuard VPN",           ORANGE,  LORANGE),
    ("102", "vm-dc",         "192.168.1.222", "Active Directory + DNS",  NAVY,    LBLUE),
    ("105", "vm-nas",        "192.168.1.219", "Fichiers CAO (Samba)",    BLUE,    LBLUE),
    ("104", "vm-erp",        "192.168.1.221", "Odoo 17 + PostgreSQL",    GREEN,   LGREEN),
    ("107", "vm-web",        "192.168.1.223", "WordPress + WooCommerce", PURPLE,  LPURPLE),
    ("103", "vm-supervision","192.168.1.224", "Prometheus + Grafana + Loki", RGBColor(0xD8,0x3B,0x01), LORANGE),
    ("101", "vm-client",     "192.168.1.211", "Poste de test",           GRAY,    LGRAY),
    ("106", "vm-clone",      "—",            "Template de base",        GRAY,    LGRAY),
]
cols = 4
cw = Inches(3.1)
ch = Inches(1.4)
gx = Inches(0.3)
gy = Inches(1.45)
for i, (vmid, name, ip, role, accent, bg_light) in enumerate(vms):
    cx = gx + (i % cols) * (cw + Inches(0.1))
    cy = gy + (i // cols) * (ch + Inches(0.15))
    rect(s, cx, cy, cw, ch, bg_light, line_color=accent, line_width=1)
    rect(s, cx, cy, Inches(0.45), ch, accent)
    txt(s, vmid, cx + Inches(0.03), cy + Inches(0.45), Inches(0.4), Inches(0.5),
        13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, cx + Inches(0.55), cy + Inches(0.1), cw - Inches(0.65), Inches(0.45),
        14, bold=True, color=NAVY)
    txt(s, ip, cx + Inches(0.55), cy + Inches(0.5), cw - Inches(0.65), Inches(0.35),
        12, color=GRAY)
    txt(s, role, cx + Inches(0.55), cy + Inches(0.85), cw - Inches(0.65), Inches(0.45),
        11, color=DARK, italic=True)

# ── SLIDE 6 — Dimensionnement ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Dimensionnement des VMs")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

rows = [
    ("vm-dc",          "2 vCPU", "4 Go",   "60 Go",      "Reco. Microsoft × 4",           BLUE),
    ("vm-nas",         "2 vCPU", "4 Go",   "40+500 Go",  "Volumétrie CAO 4 To (limite lab)", ORANGE),
    ("vm-erp",         "4 vCPU", "8 Go",   "60+100 Go",  "Lenteurs déclarées → × 2 Odoo", GREEN),
    ("vm-web",         "2 vCPU", "4 Go",   "40 Go",      "Pics campagnes promo",          PURPLE),
    ("vm-supervision", "2 vCPU", "4 Go",   "40 Go",      "Rétention logs + Grafana",      RGBColor(0xD8,0x3B,0x01)),
    ("ct-vpn",         "1 vCPU", "512 Mo", "8 Go",       "WireGuard — CT léger",          NAVY),
]
headers = ["VM", "CPU", "RAM", "Disque", "Justification"]
hw = [Inches(2.0), Inches(1.3), Inches(1.3), Inches(1.8), Inches(5.5)]
hx = Inches(0.4)
hy = Inches(1.45)
rh = Inches(0.52)
# header
for i, (h, w) in enumerate(zip(headers, hw)):
    rect(s, hx, hy, w, rh, NAVY)
    txt(s, h, hx + Inches(0.1), hy + Inches(0.1), w - Inches(0.2), rh - Inches(0.1),
        13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hx += w
for ri, (name, cpu, ram, disk, just, acc) in enumerate(rows):
    bg = WHITE if ri % 2 == 0 else RGBColor(0xF0, 0xF4, 0xFA)
    hx = Inches(0.4)
    ry = hy + rh * (ri + 1)
    vals = [name, cpu, ram, disk, just]
    for ci, (v, w) in enumerate(zip(vals, hw)):
        rect(s, hx, ry, w, rh, bg)
        if ci == 0:
            rect(s, hx, ry, Inches(0.06), rh, acc)
        txt(s, v, hx + Inches(0.12), ry + Inches(0.1), w - Inches(0.22), rh - Inches(0.1),
            12, color=NAVY if ci == 0 else DARK, bold=(ci == 0), align=PP_ALIGN.CENTER)
        hx += w

# ── SLIDE 7 — Réseau ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Réseau", "192.168.1.0/24 — réseau plat (lab) · architecture VLAN en production")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

# Schéma simplifié
rect(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(4.8), WHITE, line_color=LGRAY, line_width=0.5)
txt(s, "Hôte Proxmox VE  —  192.168.1.0/24  —  GW : 192.168.1.254",
    Inches(0.6), Inches(1.6), Inches(12.0), Inches(0.5), 14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
rect(s, Inches(0.6), Inches(2.05), Inches(12.1), Inches(0.04), LGRAY)

boxes = [
    ("ct-vpn\n.208", ORANGE),
    ("vm-dc\n.222", NAVY),
    ("vm-nas\n.219", BLUE),
    ("vm-erp\n.221", GREEN),
    ("vm-web\n.223", PURPLE),
    ("vm-sup\n.224", RGBColor(0xD8,0x3B,0x01)),
]
bx = Inches(0.7)
for label, color in boxes:
    rect(s, bx, Inches(2.2), Inches(1.9), Inches(1.5), color)
    txt(s, label, bx, Inches(2.3), Inches(1.9), Inches(1.4), 14,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bx += Inches(2.05)

txt(s, "Internet  →  Endpoint 51821 UDP  →  ct-vpn  →  réseau interne",
    Inches(0.6), Inches(4.0), Inches(12.0), Inches(0.5), 14, color=ORANGE,
    bold=True, align=PP_ALIGN.CENTER)
txt(s, "En production : IP publique METALIS + port forwarding sur le routeur",
    Inches(0.6), Inches(4.55), Inches(12.0), Inches(0.5), 13, color=GRAY,
    italic=True, align=PP_ALIGN.CENTER)

# ── SLIDE 8 — Continuité ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Continuité de service — RTO / RPO", "Fenêtre de tolérance : 20h – 4h")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

rto_rows = [
    ("Fichiers CAO",   "vm-nas", "30 min",  "4h",     "Toutes les 4h",  ORANGE),
    ("ERP Odoo",       "vm-erp", "2h",      "4h",     "Toutes les 4h",  GREEN),
    ("Site e-commerce","vm-web", "4h",      "4h",     "Toutes les 4h",  PURPLE),
    ("Active Directory","vm-dc", "1h",      "7 jours","Hebdomadaire",   NAVY),
    ("VPN",            "ct-vpn", "4h",      "N/A",    "Hebdomadaire",   ORANGE),
]
hdrs = ["Service", "VM", "RTO", "RPO", "Snapshot"]
hw2 = [Inches(3.0), Inches(2.0), Inches(1.6), Inches(1.8), Inches(3.8)]
hx = Inches(0.4); hy = Inches(1.45); rh = Inches(0.58)
for h, w in zip(hdrs, hw2):
    rect(s, hx, hy, w, rh, NAVY)
    txt(s, h, hx+Inches(0.08), hy+Inches(0.1), w-Inches(0.16), rh-Inches(0.1),
        13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hx += w
for ri, (svc, vm, rto, rpo, snap, acc) in enumerate(rto_rows):
    bg = WHITE if ri % 2 == 0 else RGBColor(0xF0,0xF4,0xFA)
    hx = Inches(0.4); ry = hy + rh*(ri+1)
    for ci, (v, w) in enumerate(zip([svc, vm, rto, rpo, snap], hw2)):
        rect(s, hx, ry, w, rh, bg)
        if ci == 0:
            rect(s, hx, ry, Inches(0.06), rh, acc)
        c = NAVY if ci == 0 else (DARK if ci != 2 else GREEN if rto in ["30 min","1h"] else ORANGE)
        txt(s, v, hx+Inches(0.12), ry+Inches(0.1), w-Inches(0.22), rh-Inches(0.1),
            12, bold=(ci==0), color=c, align=PP_ALIGN.CENTER)
        hx += w

# ── SLIDE 9 — Supervision ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Supervision", "vm-supervision — 192.168.1.224")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

sup_cards = [
    ("Prometheus",  ["Collecte métriques", "CPU · RAM · Disque · Réseau"], RGBColor(0xE6,0x52,0x1F)),
    ("Grafana",     ["Dashboards temps réel", "Visualisation centralisée"], RGBColor(0xF4,0x6C,0x00)),
    ("Loki",        ["Agrégation des logs", "Toutes les VMs"],             BLUE),
    ("Alertmanager",["Seuils configurés", "Notifications → Telegram"],     NAVY),
]
cx = Inches(0.4)
for title, lines, accent in sup_cards:
    rect(s, cx, Inches(1.5), Inches(2.9), Inches(3.0), WHITE, line_color=LGRAY, line_width=0.5)
    rect(s, cx, Inches(1.5), Inches(2.9), Inches(0.07), accent)
    txt(s, title, cx+Inches(0.15), Inches(1.65), Inches(2.6), Inches(0.5),
        18, bold=True, color=accent)
    y2 = Inches(2.25)
    for line in lines:
        txt(s, "•  "+line, cx+Inches(0.15), y2, Inches(2.6), Inches(0.4), 13, color=GRAY)
        y2 += Inches(0.45)
    cx += Inches(3.05)

txt(s, "→  Alertes automatiques sur Telegram dès qu'un seuil est dépassé",
    Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6), 16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── SLIDE 10 — Limites ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Limites et compromis")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

limits = [
    ("Nœud unique",           "Panne matérielle = tout indisponible",  "Cluster 2 nœuds + HA",         ORANGE),
    ("VPN réseau interne",    "Pas de contrôle NAT école",             "Port forwarding UDP 51821",     BLUE),
    ("500 Go CAO",            "Limite du mini-PC de lab (réel : 4 To)","6 To+ en production",           NAVY),
    ("Sauvegarde locale",     "Budget limité",                         "PBS + stockage objet S3",       GREEN),
]
y = Inches(1.5)
for title, limit, prod, acc in limits:
    rect(s, Inches(0.4), y, Inches(12.5), Inches(1.1), WHITE, line_color=LGRAY, line_width=0.5)
    rect(s, Inches(0.4), y, Inches(0.08), Inches(1.1), acc)
    txt(s, title, Inches(0.65), y+Inches(0.1), Inches(2.8), Inches(0.45), 14, bold=True, color=NAVY)
    txt(s, limit,  Inches(3.7),  y+Inches(0.1), Inches(4.0), Inches(0.45), 13, color=GRAY)
    txt(s, "→  "+prod, Inches(8.0), y+Inches(0.1), Inches(4.7), Inches(0.45), 13, bold=True, color=acc)
    txt(s, "Constat", Inches(0.65), y+Inches(0.6), Inches(2.8), Inches(0.35), 10, color=GRAY, italic=True)
    txt(s, "Limite", Inches(3.7),   y+Inches(0.6), Inches(4.0), Inches(0.35), 10, color=GRAY, italic=True)
    txt(s, "En production", Inches(8.0), y+Inches(0.6), Inches(4.7), Inches(0.35), 10, color=GRAY, italic=True)
    y += Inches(1.2)

# ── SLIDE 11 — Bilan ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Bilan — Services déployés")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

bilan = [
    ("Active Directory",       "Utilisateurs créés · postes joints au domaine", NAVY),
    ("Samba (vm-nas)",         "Partages CAO conformes aux droits AD",          BLUE),
    ("Odoo 17 (vm-erp)",       "Instance dédiée 4 vCPU / 8 Go",                GREEN),
    ("WordPress + WooCommerce","Isolé de l'ERP — API Odoo connectée",           PURPLE),
    ("WireGuard (ct-vpn)",     "Tunnel validé · config client distribuable",    ORANGE),
    ("Supervision",            "Prometheus + Grafana + Loki · alertes Telegram",RGBColor(0xD8,0x3B,0x01)),
]
cx = Inches(0.4)
cy = Inches(1.5)
for i, (title, desc, acc) in enumerate(bilan):
    x = cx + (i % 3) * Inches(4.3)
    y = cy + (i // 3) * Inches(1.8)
    rect(s, x, y, Inches(4.0), Inches(1.6), WHITE, line_color=acc, line_width=1)
    rect(s, x, y, Inches(4.0), Inches(0.07), acc)
    txt(s, "✓", x+Inches(0.12), y+Inches(0.15), Inches(0.4), Inches(0.5), 20, bold=True, color=acc)
    txt(s, title, x+Inches(0.55), y+Inches(0.15), Inches(3.3), Inches(0.5), 14, bold=True, color=NAVY)
    txt(s, desc,  x+Inches(0.15), y+Inches(0.75), Inches(3.7), Inches(0.65), 12, color=GRAY)

# ── SLIDE 12 — Améliorations ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
slide_header(s, "Améliorations envisagées", "Non réalisables en lab — architecture cible production")
rect(s, 0, Inches(1.2), W, H - Inches(1.2), OFFWHITE)

amelio = [
    ("Redondance réseau",    ["2 switches",        "Failover 4G"],                  BLUE),
    ("Double AD / DC",       ["DC secondaire",     "Réplication automatique"],      NAVY),
    ("Cluster Proxmox",      ["2 nœuds + Ceph",    "Live migration des VMs"],       GREEN),
    ("Sauvegarde externe",   ["Proxmox Backup Srv","rclone → cloud (B2/GDrive)"],   ORANGE),
]
cx = Inches(0.4)
for title, lines, acc in amelio:
    rect(s, cx, Inches(1.5), Inches(2.9), Inches(3.2), WHITE, line_color=acc, line_width=1)
    rect(s, cx, Inches(1.5), Inches(2.9), Inches(0.07), acc)
    txt(s, title, cx+Inches(0.15), Inches(1.65), Inches(2.6), Inches(0.55),
        16, bold=True, color=acc)
    y2 = Inches(2.3)
    for line in lines:
        txt(s, "—  "+line, cx+Inches(0.15), y2, Inches(2.6), Inches(0.45), 14, color=DARK)
        y2 += Inches(0.52)
    cx += Inches(3.08)

txt(s, "Objectif : éliminer tout point de défaillance unique (SPOF)",
    Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.55),
    15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── SLIDE 13 — Évolutions ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
rect(s, 0, 0, W, H, NAVY)
rect(s, 0, Inches(0), W, Inches(0.07), BLUE)
txt(s, "Évolutions envisagées", Inches(0.6), Inches(0.3), Inches(12), Inches(0.9),
    30, bold=True, color=WHITE)
rect(s, 0, Inches(1.1), W, Inches(0.04), RGBColor(0x00,0x55,0xAA))

evols = [
    ("Court terme",  "RAID logiciel (mdadm) sur vm-nas",                           ORANGE),
    ("Court terme",  "Port forwarding UDP 51821 pour VPN externe réel",             ORANGE),
    ("Moyen terme",  "WooCommerce vers hébergement cloud managé",                   BLUE),
    ("Moyen terme",  "Firewall UTM OPNsense — filtrage inter-VLAN + VPN",           BLUE),
    ("Long terme",   "Cluster Proxmox 2 nœuds — haute disponibilité",              GREEN),
    ("Long terme",   "Second site ou sync cloud pour sauvegardes (règle 3-2-1)",    GREEN),
]
y = Inches(1.35)
for term, desc, acc in evols:
    rect(s, Inches(0.5), y+Inches(0.06), Inches(1.6), Inches(0.42), acc)
    txt(s, term, Inches(0.5), y+Inches(0.08), Inches(1.6), Inches(0.4),
        11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(2.3), y, Inches(10.5), Inches(0.55), 16, color=WHITE)
    y += Inches(0.75)

txt(s, "Repo GitHub : Orealyz/metalis",
    Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
    14, color=RGBColor(0x80,0xA8,0xD4), italic=True, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"c:\Users\simon\Repos\metalis\METALIS-presentation.pptx"
prs.save(out)
print(f"Fichier généré : {out}")
